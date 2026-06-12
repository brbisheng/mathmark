"""签名注入器 - 为老师生成/注入签名

不同内容类型的注入策略:
1. PPT 文件: 用 python-pptx 修改 slide XML
2. LaTeX 文件: 字符串替换
3. Markdown/文本: 字符串替换
4. 印刷体截图: OCR + 提示
5. 手写板书: 给出建议清单

核心原则:
- **不**自动改老师原始内容(避免破坏教学逻辑)
- **提供建议**让老师在合适位置加入签名元素
- 对于 PPT/文本源, 提供可选的"自动注入"模式
"""

from __future__ import annotations

import json
import random
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Union

from ..core.types import ContentType, MathSignature, SignatureProblem
from .example_db import hash_problem

PathLike = Union[str, Path]


@dataclass
class InjectionSuggestion:
    """单条注入建议"""
    location: str  # "第 N 段", "在 ... 之后", "在 ... 之前"
    original: Optional[str]  # 原文(若为替换)
    suggestion: str  # 建议内容
    reason: str  # 为什么建议
    signature_type: str  # "marker" | "intro" | "problem" | "transition"

    def to_dict(self) -> dict:
        return {
            "location": self.location,
            "original": self.original,
            "suggestion": self.suggestion,
            "reason": self.reason,
            "signature_type": self.signature_type,
        }


@dataclass
class InjectionReport:
    """注入报告"""
    content_type: ContentType
    suggestions: List[InjectionSuggestion] = field(default_factory=list)
    automatic_injections: List[dict] = field(default_factory=list)
    summary: str = ""

    def to_dict(self) -> dict:
        return {
            "content_type": self.content_type.value,
            "suggestions": [s.to_dict() for s in self.suggestions],
            "automatic_injections": self.automatic_injections,
            "summary": self.summary,
            "n_suggestions": len(self.suggestions),
        }


# ============================================================
# 建议生成
# ============================================================

def generate_suggestions(
    signature: MathSignature,
    topic: Optional[str] = None,
    n_suggestions: int = 5,
) -> List[InjectionSuggestion]:
    """为老师生成签名注入建议

    Args:
        signature: 老师签名
        topic: 当前教学主题 (例如 "一元二次方程")
        n_suggestions: 建议数量
    """
    suggestions: List[InjectionSuggestion] = []

    # 1. 招牌例题建议
    for prob in signature.example.signature_problems[:2]:
        for variant in prob.variants[:1]:
            suggestions.append(InjectionSuggestion(
                location=f"在讲解 '{topic or '相关问题'}' 的引例部分",
                original=None,
                suggestion=f"使用招牌例题变体: {variant}",
                reason=f"这是你常用的标志性例题, 学生会记住这是你的风格",
                signature_type="problem",
            ))

    # 2. 结论标记建议
    if signature.symbol.conclusion_markers:
        marker = signature.symbol.conclusion_markers[0]
        suggestions.append(InjectionSuggestion(
            location="每个例题的最后一步",
            original=None,
            suggestion=f"用 '{marker}' 收尾(代替 '所以' 或 '故')",
            reason=f"你偏好使用 '{marker}' 作为结论标记, 这是你的风格标识",
            signature_type="marker",
        ))

    # 3. 引入语建议
    for intro in signature.step.introduction_phrases[:1]:
        suggestions.append(InjectionSuggestion(
            location="每个解题步骤的开头",
            original=None,
            suggestion=f"用 '{intro}...' 引入未知数(而非 '假设' 或 '我们设')",
            reason=f"你偏好使用 '{intro}' 作为引入语",
            signature_type="intro",
        ))

    # 4. 过渡词建议
    for trans in signature.step.transition_words[:1]:
        suggestions.append(InjectionSuggestion(
            location="代数化简后的位置",
            original=None,
            suggestion=f"加入 '{trans}' 作为过渡",
            reason=f"你偏好使用 '{trans}' 作为化简后的过渡词",
            signature_type="transition",
        ))

    # 5. 变量命名建议
    if signature.symbol.variable_primary:
        var = signature.symbol.variable_primary[0]
        suggestions.append(InjectionSuggestion(
            location="题目设问时",
            original=None,
            suggestion=f"用 '{var}' 作为主变量(而非 't' 或 'a')",
            reason=f"你偏好用 '{var}' 作为主变量名",
            signature_type="marker",
        ))

    # 限制数量
    return suggestions[:n_suggestions]


# ============================================================
# 文本注入
# ============================================================

def inject_to_text(
    text: str,
    signature: MathSignature,
    strength: float = 0.5,
) -> tuple[str, List[dict]]:
    """向纯文本注入签名元素

    注意: 这是"轻量"注入, 只做替换/添加, 不重写教学逻辑。
    返回 (modified_text, injection_log)
    """
    log: List[dict] = []
    modified = text

    # 强度控制: 0 = 不注入, 1 = 全部建议都注入
    if strength <= 0:
        return modified, log

    # 1. 替换结论标记
    common_conclusions = ["所以", "故", "因此", "得到", "可知"]
    target_marker = signature.symbol.conclusion_markers[0] if signature.symbol.conclusion_markers else "∴"

    if target_marker not in modified:
        for common in common_conclusions:
            if common in modified and random.random() < strength:
                modified = modified.replace(common, target_marker, 1)
                log.append({
                    "type": "marker_replace",
                    "from": common,
                    "to": target_marker,
                })
                break

    # 2. 在合适位置加入招牌例题(可选)
    for prob in signature.example.signature_problems[:1]:
        for variant in prob.variants[:1]:
            # 仅在 strength >= 0.5 时加入
            if strength >= 0.5 and "例题" in modified:
                # 不直接修改, 只记录建议
                log.append({
                    "type": "suggestion",
                    "suggestion": f"考虑加入变体例题: {variant}",
                })
                break

    return modified, log


def inject_to_latex(
    latex: str,
    signature: MathSignature,
    strength: float = 0.5,
) -> tuple[str, List[dict]]:
    """向 LaTeX 注入签名元素"""
    log: List[dict] = []
    modified = latex

    # LaTeX 风格的结论符注入
    latex_conclusions = [r"\\therefore", r"\\because"]
    target = signature.symbol.conclusion_markers[0] if signature.symbol.conclusion_markers else "∴"

    if target in ("∴", "∵"):
        latex_target = r"\\therefore" if target == "∴" else r"\\because"
        for common in latex_conclusions:
            if common in modified and random.random() < strength:
                # 不替换, 改用注释添加
                log.append({
                    "type": "latex_marker",
                    "suggestion": f"使用 ${latex_target}$ 渲染 '{target}'",
                })

    return modified, log


# ============================================================
# PPT 注入
# ============================================================

def inject_to_pptx(
    pptx_path: PathLike,
    signature: MathSignature,
    output_path: Optional[PathLike] = None,
    strength: float = 0.5,
) -> tuple[Path, List[dict]]:
    """向 PPT 文件注入签名元素

    实际修改的内容:
    1. 在每张幻灯片底部加入版权标记
    2. 在合适位置插入结论符提示(可选)
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches
    except ImportError:
        raise RuntimeError("python-pptx not installed. pip install python-pptx")

    log: List[dict] = []
    input_path = Path(pptx_path)
    out_path = Path(output_path) if output_path else input_path.parent / f"{input_path.stem}_signed{input_path.suffix}"

    prs = Presentation(str(input_path))

    for slide_idx, slide in enumerate(prs.slides):
        # 在每张幻灯片添加一个底部水印文本框
        try:
            left = Inches(0.2)
            top = Inches(prs.slide_height / 914400 - 0.4)  # 转换 EMU 到 inches
            width = Inches(prs.slide_width / 914400 - 0.4)
            height = Inches(0.3)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = f"© {signature.teacher_name or signature.teacher_id} | mathmark"

            # 设置样式
            p = tf.paragraphs[0]
            for run in p.runs:
                run.font.size = Pt(8)
                run.font.color.rgb = None  # 默认色

            log.append({
                "type": "pptx_watermark",
                "slide": slide_idx + 1,
                "text": f"© {signature.teacher_name or signature.teacher_id}",
            })
        except Exception as e:
            log.append({
                "type": "pptx_error",
                "slide": slide_idx + 1,
                "error": str(e),
            })

    prs.save(str(out_path))
    return out_path, log


# ============================================================
# 报告生成
# ============================================================

def create_injection_report(
    content_type: ContentType,
    signature: MathSignature,
    topic: Optional[str] = None,
    auto_inject_log: Optional[List[dict]] = None,
) -> InjectionReport:
    """创建注入报告(供 GUI 显示给老师)"""
    suggestions = generate_suggestions(signature, topic)
    summary = (
        f"为 {signature.teacher_name or signature.teacher_id} 生成了 {len(suggestions)} 条签名注入建议。\n"
        f"内容类型: {content_type.value}\n"
        f"建议: 老师在备课/发布时参考这些建议, 在合适位置加入您的标志性元素。"
    )
    return InjectionReport(
        content_type=content_type,
        suggestions=suggestions,
        automatic_injections=auto_inject_log or [],
        summary=summary,
    )


def save_injection_report(report: InjectionReport, path: PathLike) -> None:
    """保存报告为 JSON"""
    Path(path).write_text(
        json.dumps(report.to_dict(), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
